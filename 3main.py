import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# --- 페이지 설정 ---
st.set_page_config(
    page_title="📊 디저트 앱 코드 해석 PPT",
    layout="centered"
)

st.title("📊 디저트 앱 코드 해석 PPT 만들기")
st.markdown("버튼을 눌러 PPT를 다운로드하세요. 🎉")

# --- 슬라이드 내용 ---
slides = [
    {"title": "🍰 디저트 앱 소개",
     "content": ["홈베이킹 레시피 앱", "사용자가 디저트를 선택하면 재료, 레시피, 팁, 영상 제공"]},
    {"title": "🛠 페이지 설정",
     "content": ["st.set_page_config(): 앱 제목, 아이콘, 레이아웃 설정"]},
    {"title": "🎨 스타일링",
     "content": ["st.markdown('<style>...</style>'): 글자색, 배경색, 폰트, 버튼/선택창 디자인"]},
    {"title": "📌 제목과 안내문",
     "content": ["st.title(), st.markdown(): 앱 상단 제목과 안내문 표시, 사용법 안내"]},
    {"title": "🖱 디저트 선택 메뉴",
     "content": ["st.selectbox(): 사용자 선택 드롭다운 메뉴 생성"]},
    {"title": "📦 데이터 관리",
     "content": ["dict (딕셔너리): 디저트별 재료, 레시피, 팁, 영상 URL 저장", "코드 반복 없이 효율적 관리"]},
    {"title": "🔁 반복 출력",
     "content": ["for 문: 선택된 디저트 재료, 레시피, 팁 순차 출력", "화면 구성이 깔끔해짐"]},
    {"title": "🎬 영상 삽입",
     "content": ["st.video(): 유튜브 영상 삽입, 만드는 방법 영상 확인 가능"]},
    {"title": "📚 요약",
     "content": ["앱 구조: 페이지 설정 → 스타일 → 제목 → 선택 메뉴 → 데이터 → 반복 출력 → 영상", 
                 "딕셔너리와 반복문 활용으로 효율적 코드 구현"]}
]

# --- PPT 생성 함수 ---
def create_ppt(slides):
    prs = Presentation()
    bg_colors = [
        RGBColor(255, 240, 220), RGBColor(255, 230, 200), RGBColor(255, 245, 230),
        RGBColor(240, 255, 240), RGBColor(230, 240, 255), RGBColor(245, 230, 255),
        RGBColor(255, 255, 230), RGBColor(230, 255, 255), RGBColor(255, 240, 245)
    ]
    
    for i, slide_info in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_info["title"]
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)

        content_box = slide.placeholders[1]
        content_box.text = "\n".join(slide_info["content"])
        for paragraph in content_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(22)
                run.font.color.rgb = RGBColor(30, 30, 30)
        
        # 배경 색 지정
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_colors[i % len(bg_colors)]
        
    return prs

# --- 다운로드 버튼 ---
if st.button("💾 PPT 다운로드"):
    prs = create_ppt(slides)
    prs.save("디저트_앱_코드_해석.pptx")
    with open("디저트_앱_코드_해석.pptx", "rb") as f:
        st.download_button("다운로드", f, file_name="디저트_앱_코드_해석.pptx")
