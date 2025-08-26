import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt

# --- 앱 제목 ---
st.title("📊 스트림릿 앱 코드 분석 PPT 생성")

# --- PPT 생성 함수 ---
def create_ppt():
    prs = Presentation()

    # 슬라이드 1: 제목 슬라이드
    slide_layout = prs.slide_layouts[0]  # 제목 + 부제
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "🍞 디저트 추천 앱 분석"
    slide.placeholders[1].text = "스트림릿 앱 코드 구조와 기능 요약"

    # 슬라이드 2: 페이지 설정
    slide_layout = prs.slide_layouts[1]  # 제목 + 내용
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "1️⃣ 페이지 설정"
    slide.placeholders[1].text = (
        "- 앱 제목, 아이콘, 레이아웃 설정\n"
        "예시 코드:\n"
        "st.set_page_config(\n"
        "    page_title='🍞 오늘의 디저트 추천 🍰',\n"
        "    page_icon='🥐', layout='centered'\n"
        ")"
    )

    # 슬라이드 3: CSS 스타일
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "2️⃣ CSS 스타일"
    slide.placeholders[1].text = (
        "- 앱 배경, 글씨색, 버튼, 선택박스 디자인 적용\n"
        "예시 코드:\n"
        "st.markdown('<style>body {background-color:#f5e0c3; color:black;}</style>', unsafe_allow_html=True)"
    )

    # 슬라이드 4: 딕셔너리 구조
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "3️⃣ 딕셔너리 활용"
    slide.placeholders[1].text = (
        "- 디저트별 정보 저장: 재료, 레시피, 팁, 영상 링크\n"
        "예시 구조:\n"
        "dessert_info = {\n"
        "    '초코소라빵': {'ingredients':[...], 'recipe':[...], 'tips':[...], 'video': '...'},\n"
        "    '소금빵': {...}\n"
        "}"
    )

    # 슬라이드 5: 디저트 선택 & 출력
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "4️⃣ 디저트 선택 & 출력"
    slide.placeholders[1].text = (
        "- st.selectbox()로 선택\n"
        "- for문으로 재료, 레시피, 팁 출력\n"
        "- st.video()로 유튜브 영상 표시\n"
        "예시:\n"
        "for ing in info['ingredients']:\n"
        "    st.write(f'- {ing}')"
    )

    # 슬라이드 6: 주요 코드 포인트
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "5️⃣ 주요 코드 포인트"
    slide.placeholders[1].text = (
        "1. 스트림릿 기본 함수 사용: st.title, st.markdown, st.selectbox, st.write, st.video\n"
        "2. 딕셔너리로 레시피 데이터 구조화\n"
        "3. for문 반복 처리로 정보 출력\n"
        "4. CSS로 디자인 커스터마이징"
    )

    return prs

# --- PPT 다운로드 버튼 ---
if st.button("📥 발표용 PPT 만들기"):
    prs = create_ppt()
    prs.save("디저트앱_코드분석.pptx")
    st.success("✅ PPT 파일 생성 완료!")
    st.markdown("[다운로드 PPT](디저트앱_코드분석.pptx)")
